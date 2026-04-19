from flask import Flask, request
from pptx import Presentation
import os

app = Flask(__name__)

# 👉 Create uploads folder
if not os.path.exists("uploads"):
    os.makedirs("uploads")

# 👉 Analyze Text
def analyze_text(text):
    text = text.lower()
    result = ""

    if "problem" in text:
        result += "Problem Found\n"
    else:
        result += "Problem Not Found\n"

    if "solution" in text:
        result += "Solution Found\n"
    else:
        result += "Solution Not Found\n"

    if "market" in text:
        result += "Market Info Found\n"
    else:
        result += "Market Info Not Found\n"

    if "team" in text:
        result += "Team Info Found\n"
    else:
        result += "Team Info Not Found\n"

    if "revenue" in text:
        result += "Revenue Info Found\n"
    else:
        result += "Revenue Info Not Found\n"

    if "competitor" in text:
        result += "Competitor Info Found\n"
    else:
        result += "Competitor Info Not Found\n"

    return result


# 👉 Decision + Score
def get_decision(analysis):
    score = 0

    if "Problem Found" in analysis:
        score += 1
    if "Solution Found" in analysis:
        score += 1
    if "Market Info Found" in analysis:
        score += 1
    if "Team Info Found" in analysis:
        score += 1
    if "Revenue Info Found" in analysis:
        score += 1
    if "Competitor Info Found" in analysis:
        score += 1

    if score >= 5:
        decision = "Strong Buy 🚀"
    elif score >= 3:
        decision = "Hold ⚖️"
    else:
        decision = "Pass ❌"

    return score, decision


# 👉 Home Page
@app.route('/')
def home():
    return '''
    <h1 style="color:blue;">Startup Pitch Deck Analyzer</h1>

    <form action="/upload" method="post" enctype="multipart/form-data">
        <input type="file" name="file" required>
        <br><br>
        <button style="background-color:green; color:white; padding:10px;">
        Upload File
        </button>
    </form>
    '''


# 👉 Upload + Process
@app.route('/upload', methods=['POST'])
def upload():
    file = request.files['file']

    # Error handling
    if file.filename == "":
        return "❌ No file selected"

    if not file.filename.endswith(".pptx"):
        return "❌ Please upload only .pptx file"

    filepath = os.path.join("uploads", file.filename)
    file.save(filepath)

    # Extract text
    prs = Presentation(filepath)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    # Analyze
    analysis = analyze_text(text)

    # Score + Decision
    score, decision = get_decision(analysis)

    # Output UI
    return f"""
<h2 style="color:green;">Analysis Result</h2>

<pre style="background:#f4f4f4; padding:10px;">
{analysis}
</pre>

<h3>Score: {score} / 6</h3>

<h2 style="color:blue;">Investment Decision: {decision}</h2>

<br>
<a href="/">Upload Another File</a>
"""


# 👉 Run App
if __name__ == '__main__':
    app.run(debug=True)